import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import akshare as ak
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from matplotlib import rcParams 

# ======================
# 参数配置
# ======================
PRICE_DATA_PATH = "data/prices.csv"
STRATEGY_NET_PATH = "strategy_net.csv"
PARAM_RESULT_PATH = "param_search_results.csv"
OUTPUT_PPT = "quant_strategy_ppt1.pptx"
# 设置全局中文字体
rcParams['font.sans-serif'] = ['Microsoft YaHei']  # 微软雅黑
rcParams['axes.unicode_minus'] = False  # 解决负号显示问题
# ======================
# 工具函数
# ======================
def compute_metrics(series):
    """计算总收益、最大回撤、夏普比率"""
    values = series.values
    returns = values[1:] / values[:-1] - 1
    total_return = values[-1]/values[0]-1
    max_dd = (values / values.max() - 1).min()
    sharpe = returns.mean() / returns.std() * (252**0.5) if returns.std() !=0 else 0
    return total_return, max_dd, sharpe

def add_cover_slide(
    prs,
    title,
    subtitle,
    author,
    date,
):
    """
    专用封面页（研究型 / 量化答辩风格）
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白页

    # ======================
    # 主标题
    # ======================
    title_box = slide.shapes.add_textbox(
        Inches(1.0), Inches(2.2), Inches(8.0), Inches(1.5)
    )
    tf = title_box.text_frame
    tf.clear()

    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)  # 深蓝色
    p.alignment = 1  # 居中

    # ======================
    # 副标题
    # ======================
    p = tf.add_paragraph()
    p.text = subtitle
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(80, 80, 80)
    p.alignment = 1
    p.space_before = Pt(12)

    # ======================
    # 作者 & 日期
    # ======================
    info_box = slide.shapes.add_textbox(
        Inches(3.0), Inches(4.2), Inches(4.0), Inches(1.2)
    )
    tf_info = info_box.text_frame
    tf_info.clear()

    p = tf_info.paragraphs[0]
    p.text = f"作者：{author}"
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(60, 60, 60)
    p.alignment = 1

    p = tf_info.add_paragraph()
    p.text = f"日期：{date}"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(120, 120, 120)
    p.alignment = 1
def add_chart_image_slide(prs, title, fig):
    """将 matplotlib 图像嵌入 PPT 页"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    title_tf = title_box.text_frame
    title_tf.text = title
    title_tf.paragraphs[0].font.size = Pt(32)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(0,0,128)

    img_stream = BytesIO()
    fig.savefig(img_stream, format='png', bbox_inches='tight')
    plt.close(fig)
    img_stream.seek(0)
    slide.shapes.add_picture(img_stream, Inches(0.5), Inches(1.2), width=Inches(9), height=Inches(5))
    return slide

# ======================
# 特殊页：研究背景丰富化页
# ======================
def add_research_background_slide(prs, price_data_path, strategy_signal_path, strategy_net_path):
    """
    研究背景幻灯片（真实数据版本）
    - 右上：股价波动图
    - 左下：策略触发示意（实际信号）
    - 右下：小表格显示策略指标
    """
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(5))
    title_tf = title_box.text_frame
    title_tf.text = "研究背景"
    title_tf.paragraphs[0].font.size = Pt(32)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    # 左侧文字要点
    text_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4), Inches(3))
    tf = text_box.text_frame
    tf.word_wrap = True
    points = [
        "• 市场波动大，固定仓位组合难以应对短期涨跌",
        "• 固定仓位可能错失收益或承受过高回撤",
        "• 本策略基于价格涨跌触发动态加减仓",
        "• 通过动态仓位管理，控制风险、提升收益"
    ]
    for p in points:
        para = tf.add_paragraph()
        para.text = p
        para.font.size = Pt(22)
        para.font.color.rgb = RGBColor(50,50,50)

    # ---------- 右上股价波动图 ----------
    df = pd.read_csv(price_data_path, parse_dates=["date"])
    stock_code = df["stock_code"].iloc[0]  # 取第一只股票
    stock = df[df["stock_code"]==stock_code].sort_values("date")
    fig, ax = plt.subplots(figsize=(5,3))
    ax.plot(stock["date"], stock["close"], label=f"{stock_code} 收盘价", color="#1f77b4")
    ax.set_title("股票价格波动示意", fontsize=12)
    ax.set_xticks([])
    ax.set_yticks([])
    plt.tight_layout()
    img_stream = BytesIO()
    fig.savefig(img_stream, format='png', bbox_inches='tight', dpi=200)
    plt.close(fig)
    img_stream.seek(0)
    slide.shapes.add_picture(img_stream, Inches(5), Inches(1.2), width=Inches(4.5), height=Inches(3))

    # ---------- 左下策略触发示意 ----------
    df_signal = pd.read_csv(strategy_signal_path, parse_dates=["date"]).sort_values("date")
    # 策略信号：1=加仓, -1=减仓, 0=不操作
    signal_dates = df_signal["date"]
    signal_values = df_signal["signal"]

    fig, ax = plt.subplots(figsize=(5,2))
    ax.plot(signal_dates, signal_values, marker='o', color="#1f77b4")
    for d, s in zip(signal_dates, signal_values):
        if s > 0:
            ax.text(d, s, "+", fontsize=10, color="#1f77b4")
        elif s < 0:
            ax.text(d, s, "-", fontsize=10, color="#1f77b4")
    ax.set_xticks([])
    ax.set_yticks([])
    ax.set_title("策略触发示意（实际信号）", fontsize=12)
    plt.tight_layout()
    img_stream2 = BytesIO()
    fig.savefig(img_stream2, format='png', bbox_inches='tight', dpi=200)
    plt.close(fig)
    img_stream2.seek(0)
    slide.shapes.add_picture(img_stream2, Inches(0.5), Inches(4.3), width=Inches(4.5), height=Inches(2))

    # ---------- 右下小表格 ----------
    df_net = pd.read_csv(strategy_net_path, parse_dates=["date"]).set_index("date")
    total_return = df_net["net_value"].iloc[-1]/df_net["net_value"].iloc[0]-1
    max_drawdown = (df_net["net_value"]/df_net["net_value"].cummax()-1).min()

    rows, cols = 3, 2
    table = slide.shapes.add_table(rows, cols, Inches(5), Inches(4.3), Inches(4), Inches(2)).table
    table.cell(0,0).text = "示例指标"
    table.cell(0,1).text = "数值"
    table.cell(1,0).text = "总收益率"
    table.cell(1,1).text = f"{total_return*100:.2f}%"
    table.cell(2,0).text = "最大回撤"
    table.cell(2,1).text = f"{max_drawdown*100:.2f}%"
    for r in range(rows):
        for c in range(cols):
            table.cell(r,c).text_frame.paragraphs[0].font.size = Pt(14)

    return slide

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from io import BytesIO
from PIL import Image

def add_data_source_slide_with_figures(prs, price_hist_path, vol_hist_path):
    """
    数据来源页 + 分析图表
    price_hist_path: 股票均价分布直方图路径
    vol_hist_path: 股票波动率直方图路径
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白页

    # --------------------
    # 左侧文字说明
    # --------------------
    tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(4.5), Inches(5))
    tf = tx_box.text_frame
    tf.word_wrap = True

    p = tf.add_paragraph()
    p.text = "数据来源与探索分析"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    p = tf.add_paragraph()
    p.text = (
        "• A股前 N 只股票历史日线数据（2022–2023）\n"
        "• 数据通过 AkShare 获取\n"
        "• 保证每只股票交易日一致（交易日交集筛选）\n"
        "• 对股票价格与波动性进行了探索性分析"
    )
    p.font.size = Pt(22)
    p.space_after = Pt(10)

    # --------------------
    # 右侧插入图表
    # --------------------
    # 股票均价直方图
    slide.shapes.add_picture(price_hist_path, Inches(5), Inches(0.5), width=Inches(4), height=Inches(2.5))

    # 波动率直方图
    slide.shapes.add_picture(vol_hist_path, Inches(5), Inches(3), width=Inches(4), height=Inches(2.5))

    return slide

# ======================
# 插入策略设计 + 参数学习页（整体上移 1 英寸）
# ======================
def insert_strategy_param_slide(prs, insert_idx=3):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白页

    # ---------- 左侧文字 ----------
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(4.5), Inches(5))  # 上移 0.3 -> 0.2
    tf = left_box.text_frame
    tf.word_wrap = True

    p = tf.add_paragraph()
    p.text = "策略设计与参数自动学习"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    points = [
        "核心思想：价格涨跌触发动态加减仓",
        "卖出条件：价格涨幅 ≥ 5%",
        "买入条件：价格跌幅 ≥ 5%",
        "调仓比例：每次 ±5%",
        "最大仓位：初始资金 2 倍",
        "",
        "参数自动学习方法（Parameter Learning）：",
        "通过历史数据验证与筛选关键参数，提高策略稳健性与客观性。"
    ]
    for pt in points:
        para = tf.add_paragraph()
        para.text = pt
        para.font.size = Pt(22)
        para.font.color.rgb = RGBColor(50,50,50)

    # ---------- 右侧流程图（压缩版） ----------
    fig, ax = plt.subplots(figsize=(3.5,2.5))  # 压缩高度
    ax.axis('off')
    nodes = {
        "生成参数组合": (0.5, 0.9),
        "历史回测": (0.5, 0.7),
        "计算指标": (0.5, 0.5),
        "记录结果": (0.5, 0.3),
        "选择最优参数": (0.5, 0.1)
    }
    for name, (x, y) in nodes.items():
        ax.text(x, y, name, ha='center', va='center',
                bbox=dict(boxstyle="round,pad=0.3", facecolor="#A6CEE3", edgecolor="b"))
    arrowprops=dict(arrowstyle="->", color='black', lw=1.5)
    for start, end in zip(list(nodes.values())[:-1], list(nodes.values())[1:]):
        ax.annotate("", xy=end, xytext=start, arrowprops=arrowprops)
    plt.tight_layout()
    img_stream = BytesIO()
    fig.savefig(img_stream, format='png', bbox_inches='tight', dpi=200)
    plt.close(fig)
    img_stream.seek(0)
    slide.shapes.add_picture(img_stream, Inches(5), Inches(0.8), width=Inches(4), height=Inches(2.5))  # 上移 0.2->0.8

    # ---------- 参数表格 ----------
    rows, cols = 4, 4
    table_top = Inches(5)   # 上移 1 英寸
    table_height = Inches(1)  # 压缩高度
    table = slide.shapes.add_table(rows, cols, Inches(0.5), table_top, Inches(9), table_height).table

    table.cell(0,0).text = "参数"
    table.cell(0,1).text = "含义"
    table.cell(0,2).text = "搜索范围"
    table.cell(0,3).text = "对策略影响"
    params = [
        ("PRICE_TRIGGER", "触发买卖操作的价格涨跌幅", "3%,5%,7%", "阈值小→频繁交易，阈值大→滞后交易"),
        ("TRADE_RATIO", "每次加减仓比例", "3%,5%,10%", "大→波动大，小→调仓缓慢"),
        ("MAX_POSITION_MULT", "单只股票最大仓位倍数", "1.5,2.0", "控制单股暴露风险")
    ]
    for i, row in enumerate(params, start=1):
        for j, val in enumerate(row):
            table.cell(i,j).text = val
            table.cell(i,j).text_frame.paragraphs[0].font.size = Pt(14)

    # ---------- 插入到指定位置 ----------
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[-1])
    xml_slides.insert(insert_idx, slides[-1])

def download_hs300(path="HS300.csv"):
    """
    下载沪深300指数日线历史行情，并保存为 CSV
    """
    # 用正确的 akshare 接口获取日线数据
    df = ak.stock_zh_index_daily_em(symbol="sh000300")

    # 返回的列里通常包括 date, open, high, low, close, volume 等
    # 我们只保留 date 和 close
    df = df[["date", "close"]]
    df["date"] = pd.to_datetime(df["date"])
    df.sort_values("date", inplace=True)

    df.to_csv(path, index=False)
    print(f"HS300 数据已保存到 {path}")
    return path


hs300_csv = download_hs300()

# -----------------------------
# 2️⃣ 原函数改写（加入 HS300 绘图）
# -----------------------------
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import pandas as pd
import matplotlib.pyplot as plt
from io import BytesIO

def add_backtest_vs_benchmark_slide(prs, strategy_net_path, price_data_path, hs300_path=None, test_metrics=None):
    """
    添加回测 vs 基准净值对比幻灯片
    左侧：基准说明
    右侧：策略净值 vs 等权组合 vs HS300折线图
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白页

    # ---------- 左侧基准说明 ----------
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(4.5), Inches(5))
    tf = left_box.text_frame
    tf.word_wrap = True

    p = tf.add_paragraph()
    p.text = "基准选择与说明"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    points = [
        "• 等权组合（EqualWeight）：所有选定股票等权持仓",
        "• 沪深300指数（HS300）：市场代表指数",
        "• 对齐起点日期，统一归一化比较",
        "• 对比策略净值曲线，直观评估策略超额收益"
    ]
    for pt in points:
        para = tf.add_paragraph()
        para.text = pt
        para.font.size = Pt(22)
        para.font.color.rgb = RGBColor(50,50,50)

    # ---------- 右侧回测净值折线图 ----------

    # 1️⃣ 读取策略净值
    df_strategy = pd.read_csv(strategy_net_path, parse_dates=["date"]).set_index("date").sort_index()
    df_strategy["net_value"] = df_strategy["net_value"].astype(float)

    # 2️⃣ 读取股票价格，计算等权组合净值
    df_prices = pd.read_csv(price_data_path, parse_dates=["date"]).sort_values(["stock_code","date"])
    df_prices["daily_ret"] = df_prices.groupby("stock_code")["close"].pct_change()
    ew_ret = df_prices.groupby("date")["daily_ret"].mean().dropna()
    ew_net = (1 + ew_ret).cumprod()

    # 3️⃣ 读取HS300数据
    if hs300_path:
        df_hs300 = pd.read_csv(hs300_path, parse_dates=["date"]).set_index("date").sort_index()
    else:
        df_hs300 = None

    # 4️⃣ 对齐日期
    common_dates = df_strategy.index.intersection(ew_net.index)
    if df_hs300 is not None:
        common_dates = common_dates.intersection(df_hs300.index)

    # 5️⃣ 策略净值 & 等权组合归一化
    df_strategy = df_strategy.loc[common_dates]
    df_strategy["net_value"] /= df_strategy["net_value"].iloc[0]

    ew_net = ew_net.loc[common_dates]
    ew_net /= ew_net.iloc[0]

    if df_hs300 is not None:
        hs300_net = df_hs300.loc[common_dates, "close"]
        hs300_net = hs300_net.astype(float)
        hs300_net /= hs300_net.iloc[0]
    else:
        hs300_net = None

    # 打印检查起点
    print("策略净值起点：", df_strategy["net_value"].iloc[0])
    print("等权组合起点：", ew_net.iloc[0])
    if hs300_net is not None:
        print("HS300起点：", hs300_net.iloc[0])

    # 6️⃣ 绘图
    fig, ax = plt.subplots(figsize=(5,4))
    ax.plot(common_dates, df_strategy["net_value"], label="策略净值", color="#1f77b4")
    ax.plot(common_dates, ew_net, label="等权组合净值", color="#ff7f0e")
    if hs300_net is not None:
        ax.plot(common_dates, hs300_net, label="沪深300", color="#2ca02c")
    ax.set_title("策略净值 vs 基准", fontsize=14)
    ax.set_xlabel("日期", fontsize=12)
    ax.set_ylabel("归一化净值", fontsize=12)
    ax.legend(fontsize=10)
    plt.xticks(rotation=30)
    plt.tight_layout()

    # 7️⃣ 保存图片到 PPT
    img_stream = BytesIO()
    fig.savefig(img_stream, format='png', bbox_inches='tight', dpi=200)
    plt.close(fig)
    img_stream.seek(0)
    slide.shapes.add_picture(img_stream, Inches(5), Inches(0.5), width=Inches(4.5), height=Inches(4))
    # ---------- 添加测试集绩效指标表格 ----------
    if test_metrics:
        rows = len(test_metrics) + 1
        cols = 2
        left = Inches(0.5)
        top = Inches(4.8)
        width = Inches(8)
        height = Inches(3)

        table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
        table = table_shape.table

        # 表头
        table.cell(0,0).text = "指标"
        table.cell(0,1).text = "数值"
        for i in range(2):
            table.cell(0,i).text_frame.paragraphs[0].font.bold = True
            table.cell(0,i).text_frame.paragraphs[0].font.size = Pt(22)
            table.cell(0,i).text_frame.paragraphs[0].font.color.rgb = RGBColor(0,51,102)

        # 填充数据
        for idx, (k,v) in enumerate(test_metrics.items()):
            table.cell(idx+1,0).text = k
            table.cell(idx+1,1).text = f"{v:.4f}"
            for j in range(2):
                table.cell(idx+1,j).text_frame.paragraphs[0].font.size = Pt(22)
                table.cell(idx+1,j).text_frame.paragraphs[0].font.color.rgb = RGBColor(50,50,50)

    return slide



def add_summary_outlook_slide(
    prs,
    conclusions,
    outlooks
):
    """
    总结与展望页（研究型 / 课程答辩风格）
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白页

    # ======================
    # 标题
    # ======================
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(9), Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "总结与展望"
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    # ======================
    # 左侧：总结
    # ======================
    left_box = slide.shapes.add_textbox(
        Inches(0.6), Inches(1.4), Inches(4.2), Inches(4.5)
    )
    tf_left = left_box.text_frame
    tf_left.word_wrap = True

    p = tf_left.paragraphs[0]
    p.text = "研究总结"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(64, 64, 64)

    for text in conclusions:
        para = tf_left.add_paragraph()
        para.text = text
        para.font.size = Pt(22)
        para.font.color.rgb = RGBColor(80, 80, 80)
        para.space_before = Pt(6)

    # ======================
    # 右侧：展望
    # ======================
    right_box = slide.shapes.add_textbox(
        Inches(5.2), Inches(1.4), Inches(4.2), Inches(4.5)
    )
    tf_right = right_box.text_frame
    tf_right.word_wrap = True

    p = tf_right.paragraphs[0]
    p.text = "未来展望"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(64, 64, 64)

    for text in outlooks:
        para = tf_right.add_paragraph()
        para.text = text
        para.font.size = Pt(22)
        para.font.color.rgb = RGBColor(80, 80, 80)
        para.space_before = Pt(6)
# ======================
# 特殊页：参数敏感性热力图（直接插入 PNG）
# ======================
def add_param_sensitivity_slide(prs, img_path=r"D:\database\stocknew\strategy\param_sensitivity_heatmap.png"):
    from pptx.util import Inches

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白页

    # 插入标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "参数敏感性分析"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    # 插入 PNG 图片（全页宽度，高度自适应）
    slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.2), width=Inches(9))

# ======================
# 主流程
# ======================
def main():
    prs = Presentation()

    # 第1页：封面
    add_cover_slide(
        prs,
        title="基于价格波动的动态仓位管理策略研究",
        subtitle="—— 量化投资策略的构建、回测与基准比较 ——",
        author="李萍",
        date="2026 年 1 月 12 日"
    )

    # 第2页：研究背景丰富化页
    add_research_background_slide(
    prs,
    price_data_path=PRICE_DATA_PATH,
    strategy_signal_path="data/strategy_signal.csv",
    strategy_net_path=STRATEGY_NET_PATH
)

    # 第3页：数据来源（图文并列）
    add_data_source_slide_with_figures(
    prs,
    price_hist_path="analysis_figures/price_mean_hist.png",
    vol_hist_path="analysis_figures/volatility_hist.png"
)

    # 第4页：策略设计 + 参数学习
    insert_strategy_param_slide(prs)

    # 第5页：参数敏感性热力图
    add_param_sensitivity_slide(prs)

    # 第6页：基准对比+回测结果（图表）
    metrics = {"Total Return": -0.0048, "Max Drawdown": -0.0249, "Sharpe Ratio": -0.0952}
    add_backtest_vs_benchmark_slide(
    prs,
    strategy_net_path=STRATEGY_NET_PATH,
    price_data_path=PRICE_DATA_PATH,
    hs300_path=hs300_csv, test_metrics=metrics
)

    # 第7页：总结与展望
    add_summary_outlook_slide(
        prs,
        conclusions=[
            "基于价格波动的动态仓位管理策略能够有效捕捉市场短期波动机会。",
            "在历史回测中，该策略相较于等权组合与指数基准表现出更优的风险收益特征。",
            "参数自动学习方法提高了策略构建的客观性与稳健性。"
        ],
        outlooks=[
            "引入多因子信号（如动量、波动率、成交量）以增强交易决策。",
            "结合风险控制机制（止损、最大回撤约束）进一步优化策略稳定性。",
            "扩展至不同时间尺度数据（高频或长期）以验证策略泛化能力。"
        ]
    )

    # 保存 PPT
    prs.save(OUTPUT_PPT)
    print(f"PPT 已生成: {OUTPUT_PPT}")

if __name__ == "__main__":
    main()
